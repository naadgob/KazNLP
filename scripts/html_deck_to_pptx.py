"""Export HTML slide deck (1920x1080) to PowerPoint via Playwright screenshots."""

from __future__ import annotations

import argparse
import tempfile
from pathlib import Path

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches

SLIDE_W = 1920
SLIDE_H = 1080


def export_deck(html_path: Path, pptx_path: Path) -> int:
    html_path = html_path.resolve()
    pptx_path = pptx_path.resolve()
    pptx_path.parent.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory() as tmp:
        tmp_dir = Path(tmp)
        shots: list[Path] = []

        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page(viewport={"width": SLIDE_W, "height": SLIDE_H})
            page.goto(html_path.as_uri(), wait_until="networkidle")
            page.wait_for_timeout(800)

            slide_count = page.evaluate(
                "() => document.querySelectorAll('.slide').length"
            )

            for i in range(slide_count):
                page.evaluate(
                    """(idx) => {
                        const slides = document.querySelectorAll('.slide');
                        const total = slides.length;
                        slides.forEach((el, j) => el.classList.toggle('active', j === idx));
                        const stage = document.getElementById('deck-stage');
                        if (stage) {
                            stage.style.transform = 'none';
                            stage.style.transformOrigin = 'top left';
                        }
                        const counter = document.getElementById('deck-counter');
                        if (counter) {
                            counter.textContent = String(idx + 1).padStart(2, '0') + ' / ' + String(total);
                        }
                    }""",
                    i,
                )
                page.wait_for_timeout(200)
                shot = tmp_dir / f"slide_{i + 1:02d}.png"
                page.locator("#deck-stage").screenshot(path=str(shot))
                shots.append(shot)

            browser.close()

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        for shot in shots:
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(
                str(shot), Inches(0), Inches(0), width=prs.slide_width
            )

        prs.core_properties.title = "KazNLP — презентация"
        prs.core_properties.author = "Bogdan Savelyev"
        prs.save(pptx_path)

    return slide_count


def main() -> None:
    parser = argparse.ArgumentParser(description="HTML deck → PowerPoint")
    parser.add_argument("html", type=Path, help="Path to HTML deck")
    parser.add_argument(
        "-o",
        "--output",
        type=Path,
        default=None,
        help="Output .pptx path",
    )
    args = parser.parse_args()

    out = args.output or args.html.with_suffix(".pptx")
    n = export_deck(args.html, out)
    print(f"Wrote {out} ({n} slides)")


if __name__ == "__main__":
    main()
